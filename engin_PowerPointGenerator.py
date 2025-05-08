from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation and add blank slide
prs = Presentation()
slide_width = prs.slide_width
slide_height = prs.slide_height
slide = prs.slides.add_slide(prs.slide_layouts[6])

# === Background ===
slide.shapes.add_picture("static/images/background.jpg", 0, 0, width=slide_width, height=slide_height)

# === Title ===
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
title_frame = title_box.text_frame
title_frame.clear()
title_run = title_frame.paragraphs[0].add_run()
title_run.text = "AFTERSHOCKS"
title_run.font.size = Pt(60)
title_run.font.bold = True
title_run.font.name = "Impact"
title_run.font.color.rgb = RGBColor(255, 69, 0)  # Red-orange

# === Tagline ===
tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(0.6))
tagline_frame = tagline_box.text_frame
tagline_run = tagline_frame.paragraphs[0].add_run()
tagline_run.text = "The Earth cracked open. So did the Gates of the Dead."
tagline_run.font.size = Pt(22)
tagline_run.font.italic = True
tagline_run.font.color.rgb = RGBColor(255, 255, 255)

# # === CAST ===
# cast = [
#     ("static\images\dwayne.jpg", "Dwayne Johnson"),
#     ("static\images\jennifer.jpg", "Jennifer Lawrence"),
#     ("static\images\trejo.jpg", "Danny Trejo"),
#     ("static\images\ziering.jpg", "Ian Ziering"),
#     ("static\images\reid.jpg", "Tara Reid"),
# ]

# import os

# base_path = "static/images"
# cast = [
#     (os.path.join(base_path, "dwayne.jpg"), "Dwayne Johnson"),
#     (os.path.join(base_path, "jennifer.jpg"), "Jennifer Lawrence"),
#     (os.path.join(base_path, "trejo.jpg"), "Danny Trejo"),
#     (os.path.join(base_path, "ziering.jpg"), "Ian Ziering"),
#     (os.path.join(base_path, "reid.jpg"), "Tara Reid"),
# ]

cast = [
    ("static/images/dwayne.jpg", "Dwayne Johnson"),
    ("static/images/jennifer.jpg", "Jennifer Lawrence"),
    ("static/images/trejo.jpg", "Danny Trejo"),
    ("static/images/ziering.jpg", "Ian Ziering"),
    ("static/images/reid.jpg", "Tara Reid"),
]



x_positions = [0.4, 2.0, 3.6, 5.2, 6.8]
y_image = 2.4
y_label = 4.0
img_width = Inches(1.4)

for i, (img_file, name) in enumerate(cast):
    x = Inches(x_positions[i])
    slide.shapes.add_picture(img_file, x, Inches(y_image), width=img_width)
    label_box = slide.shapes.add_textbox(x, Inches(y_label), img_width, Inches(0.5))
    label_frame = label_box.text_frame
    label_frame.paragraphs[0].text = name
    label_frame.paragraphs[0].font.size = Pt(12)
    label_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 0)
    label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# === Director ===
director_box = slide.shapes.add_textbox(Inches(6.4), Inches(0.4), Inches(3), Inches(0.5))
director_frame = director_box.text_frame
director_frame.paragraphs[0].text = "Directed by James DeMonaco"
director_frame.paragraphs[0].font.size = Pt(14)
director_frame.paragraphs[0].font.color.rgb = RGBColor(173, 216, 230)

# === Plot (Moved to bottom) ===
plot_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(8.5), Inches(1.2))
plot_frame = plot_box.text_frame
plot_frame.word_wrap = True
plot_paragraph = plot_frame.add_paragraph()
plot_paragraph.text = (
    "After a massive quake hits Goodsprings, something ancient awakens beneath the rubble. "
    "A team led by Dwayne Johnson must survive the undead to stop a world-ending curse."
)
plot_paragraph.font.size = Pt(14)
plot_paragraph.font.color.rgb = RGBColor(240, 240, 240)
plot_paragraph.alignment = PP_ALIGN.LEFT

# === Call to Action ===
cta_box = slide.shapes.add_textbox(Inches(2.5), Inches(6.8), Inches(5), Inches(1))
cta_frame = cta_box.text_frame
cta_run = cta_frame.paragraphs[0].add_run()
cta_run.text = "🔥 Coming Soon • Only In Theaters"
cta_run.font.size = Pt(22)
cta_run.font.bold = True
cta_run.font.color.rgb = RGBColor(255, 0, 0)
cta_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save the final presentation
prs.save("Aftershocks_Movie_Poster.pptx")
print("✅ Slide updated: Poster saved as Aftershocks_Movie_Poster.pptx")
