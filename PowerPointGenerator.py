import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import os
import sys
from PIL import Image, ImageTk
import threading

# This would be used to import your PowerPoint generation code
# from pptx_generator import generate_poster

class MoviePosterApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Movie Poster PowerPoint Generator By AJ")
        self.root.geometry("900x700")
        self.root.minsize(800, 600)
        
        # Apply a theme if available
        try:
            self.root.tk.call("source", "azure.tcl")
            self.root.tk.call("set_theme", "light")
        except:
            pass  # If theme is not available, use default

        self.setup_variables()
        self.create_widgets()
        self.setup_layout()
    
    def setup_variables(self):
        """Initialize all data variables used in the app"""
        self.title_var = tk.StringVar(value="AFTERSHOCKS")
        self.tagline_var = tk.StringVar(value="The Earth cracked open. So did the Gates of the Dead.")
        self.director_var = tk.StringVar(value="James DeMonaco")
        self.plot_var = tk.StringVar(value="After a massive quake hits Goodsprings, something ancient awakens beneath the rubble. A team led by Dwayne Johnson must survive the undead to stop a world-ending curse.")
        self.cta_var = tk.StringVar(value="🔥 Coming Soon • Only In Theaters")
        self.bg_image_path = tk.StringVar(value="static/images/background.jpg")
        
        # Cast is stored as a list of dictionaries
        self.cast = [
            {"name": "Dwayne Johnson", "image_path": "static/images/dwayne.jpg"},
            {"name": "Jennifer Lawrence", "image_path": "static/images/jennifer.jpg"},
            {"name": "Danny Trejo", "image_path": "static/images/trejo.jpg"},
            {"name": "Ian Ziering", "image_path": "static/images/ziering.jpg"},
            {"name": "Tara Reid", "image_path": "static/images/reid.jpg"},
        ]
        
        # Status variable for the progress bar
        self.status_var = tk.StringVar(value="Ready")
        self.progress_var = tk.IntVar(value=0)
        
        # Flag for saving code
        self.save_code_var = tk.BooleanVar(value=False)
        
        # Variable to store the selected save filename
        self.save_filename = None
        
    def create_widgets(self):
        """Create all widgets for the application"""
        # Create a notebook for tabs
        self.notebook = ttk.Notebook(self.root)
        
        # Create frames for each tab
        self.main_frame = ttk.Frame(self.notebook)
        self.cast_frame = ttk.Frame(self.notebook)
        self.preview_frame = ttk.Frame(self.notebook)
        
        # Add the frames to the notebook
        self.notebook.add(self.main_frame, text="Movie Details")
        self.notebook.add(self.cast_frame, text="Cast Members")
        self.notebook.add(self.preview_frame, text="Preview & Generate")
        
        # Create bottom status bar
        self.status_frame = ttk.Frame(self.root)
        self.status_label = ttk.Label(self.status_frame, textvariable=self.status_var)
        self.progress_bar = ttk.Progressbar(
            self.status_frame, 
            variable=self.progress_var, 
            maximum=100,
            length=300,
            mode="determinate"
        )
        
        # Create widgets for main tab
        self.create_main_tab_widgets()
        
        # Create widgets for cast tab
        self.create_cast_tab_widgets()
        
        # Create widgets for preview tab
        self.create_preview_tab_widgets()
        
    def create_main_tab_widgets(self):
        """Create widgets for the main movie details tab"""
        # Title
        title_label = ttk.Label(self.main_frame, text="Movie Title:")
        self.title_entry = ttk.Entry(self.main_frame, textvariable=self.title_var, width=40)
        
        # Tagline
        tagline_label = ttk.Label(self.main_frame, text="Tagline:")
        self.tagline_entry = ttk.Entry(self.main_frame, textvariable=self.tagline_var, width=60)
        
        # Director
        director_label = ttk.Label(self.main_frame, text="Director:")
        self.director_entry = ttk.Entry(self.main_frame, textvariable=self.director_var, width=40)
        
        # Plot
        plot_label = ttk.Label(self.main_frame, text="Plot:")
        self.plot_text = tk.Text(self.main_frame, height=5, width=60, wrap=tk.WORD)
        self.plot_text.insert("1.0", self.plot_var.get())
        
        # Call to Action
        cta_label = ttk.Label(self.main_frame, text="Call to Action:")
        self.cta_entry = ttk.Entry(self.main_frame, textvariable=self.cta_var, width=60)
        
        # Background Image
        bg_label = ttk.Label(self.main_frame, text="Background Image:")
        bg_frame = ttk.Frame(self.main_frame)
        self.bg_entry = ttk.Entry(bg_frame, textvariable=self.bg_image_path, width=50)
        bg_browse_btn = ttk.Button(
            bg_frame, 
            text="Browse", 
            command=lambda: self.browse_image(self.bg_image_path)
        )
        
        # Place main tab widgets
        title_label.grid(row=0, column=0, sticky="w", padx=10, pady=10)
        self.title_entry.grid(row=0, column=1, sticky="w", padx=10, pady=10)
        
        tagline_label.grid(row=1, column=0, sticky="w", padx=10, pady=10)
        self.tagline_entry.grid(row=1, column=1, sticky="w", padx=10, pady=10)
        
        director_label.grid(row=2, column=0, sticky="w", padx=10, pady=10)
        self.director_entry.grid(row=2, column=1, sticky="w", padx=10, pady=10)
        
        plot_label.grid(row=3, column=0, sticky="nw", padx=10, pady=10)
        self.plot_text.grid(row=3, column=1, sticky="w", padx=10, pady=10)
        
        cta_label.grid(row=4, column=0, sticky="w", padx=10, pady=10)
        self.cta_entry.grid(row=4, column=1, sticky="w", padx=10, pady=10)
        
        bg_label.grid(row=5, column=0, sticky="w", padx=10, pady=10)
        bg_frame.grid(row=5, column=1, sticky="w", padx=10, pady=10)
        self.bg_entry.pack(side=tk.LEFT, fill=tk.X, expand=True)
        bg_browse_btn.pack(side=tk.RIGHT, padx=5)
        
    def create_cast_tab_widgets(self):
        """Create widgets for the cast members tab"""
        # Controls for adding/removing cast members
        controls_frame = ttk.Frame(self.cast_frame)
        add_btn = ttk.Button(
            controls_frame, 
            text="Add Cast Member", 
            command=self.add_cast_member
        )
        
        # Scrollable frame for cast members
        self.cast_canvas = tk.Canvas(self.cast_frame)
        scrollbar = ttk.Scrollbar(self.cast_frame, orient="vertical", command=self.cast_canvas.yview)
        self.cast_canvas.configure(yscrollcommand=scrollbar.set)
        
        self.cast_scroll_frame = ttk.Frame(self.cast_canvas)
        self.cast_scroll_frame_id = self.cast_canvas.create_window(
            (0, 0), 
            window=self.cast_scroll_frame, 
            anchor="nw"
        )
        
        # Place cast tab widgets
        controls_frame.pack(fill=tk.X, padx=10, pady=10)
        add_btn.pack(side=tk.LEFT)
        
        self.cast_canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=10, pady=10)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # Configure the canvas
        self.cast_scroll_frame.bind(
            "<Configure>",
            lambda e: self.cast_canvas.configure(scrollregion=self.cast_canvas.bbox("all"))
        )
        self.cast_canvas.bind("<Configure>", self.on_canvas_configure)
        
        # Initialize cast member widgets
        self.cast_member_frames = []
        for i, _ in enumerate(self.cast):
            self.create_cast_member_widget(i)
        
    def on_canvas_configure(self, event):
        """Handle canvas resize"""
        self.cast_canvas.itemconfig(self.cast_scroll_frame_id, width=event.width)
    
    def create_cast_member_widget(self, index):
        """Create a widget for a single cast member"""
        # Create a frame for this cast member
        member_frame = ttk.LabelFrame(
            self.cast_scroll_frame, 
            text=f"Cast Member #{index + 1}"
        )
        
        # Name
        name_label = ttk.Label(member_frame, text="Name:")
        name_var = tk.StringVar(value=self.cast[index]["name"])
        name_entry = ttk.Entry(member_frame, textvariable=name_var, width=30)
        
        # Update name when changed
        def update_name(*args):
            self.cast[index]["name"] = name_var.get()
        name_var.trace_add("write", update_name)
        
        # Image
        image_label = ttk.Label(member_frame, text="Image:")
        image_frame = ttk.Frame(member_frame)
        image_var = tk.StringVar(value=self.cast[index]["image_path"])
        image_entry = ttk.Entry(image_frame, textvariable=image_var, width=40)
        
        # Update image when changed
        def update_image(*args):
            self.cast[index]["image_path"] = image_var.get()
        image_var.trace_add("write", update_image)
        
        image_browse_btn = ttk.Button(
            image_frame, 
            text="Browse", 
            command=lambda: self.browse_image(image_var)
        )
        
        # Remove button
        remove_btn = ttk.Button(
            member_frame, 
            text="Remove", 
            command=lambda: self.remove_cast_member(index)
        )
        
        # Place widgets in the frame
        name_label.grid(row=0, column=0, sticky="w", padx=5, pady=5)
        name_entry.grid(row=0, column=1, sticky="w", padx=5, pady=5)
        
        image_label.grid(row=1, column=0, sticky="w", padx=5, pady=5)
        image_frame.grid(row=1, column=1, sticky="w", padx=5, pady=5)
        image_entry.pack(side=tk.LEFT, fill=tk.X, expand=True)
        image_browse_btn.pack(side=tk.RIGHT, padx=5)
        
        remove_btn.grid(row=2, column=1, sticky="e", padx=5, pady=5)
        
        # Add the frame to our list
        member_frame.pack(fill=tk.X, padx=10, pady=5)
        self.cast_member_frames.append({
            "frame": member_frame,
            "name_var": name_var,
            "image_var": image_var
        })
    
    def create_preview_tab_widgets(self):
        """Create widgets for the preview & generate tab"""
        # Preview pane (would hold a mock image)
        preview_label = ttk.Label(self.preview_frame, text="Preview:")
        self.preview_image_label = ttk.Label(self.preview_frame)
        
        # Generate button
        generate_frame = ttk.Frame(self.preview_frame)
        save_code_check = ttk.Checkbutton(
            generate_frame,
            text="Save Python Code",
            variable=self.save_code_var
        )
        generate_btn = ttk.Button(
            generate_frame, 
            text="Generate PowerPoint",
            command=self.generate_powerpoint
        )
        
        # Place preview tab widgets
        preview_label.pack(anchor="w", padx=10, pady=5)
        self.preview_image_label.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        generate_frame.pack(fill=tk.X, padx=10, pady=10)
        save_code_check.pack(side=tk.LEFT, padx=5)
        generate_btn.pack(side=tk.RIGHT, padx=5)
        
        # Set a mock preview image
        self.update_preview()
    
    def setup_layout(self):
        """Set up the overall layout of the application"""
        self.notebook.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        self.status_frame.pack(fill=tk.X, side=tk.BOTTOM, padx=10, pady=5)
        self.status_label.pack(side=tk.LEFT, padx=5)
        self.progress_bar.pack(side=tk.RIGHT, padx=5)
    
    def browse_image(self, string_var):
        """Opens a file dialog to select an image"""
        filename = filedialog.askopenfilename(
            title="Select an image",
            filetypes=(
                ("Image files", "*.jpg *.jpeg *.png *.gif"),
                ("All files", "*.*")
            )
        )
        if filename:
            string_var.set(filename)
            # Update preview if we've changed the background
            if string_var == self.bg_image_path:
                self.update_preview()
    
    def add_cast_member(self):
        """Add a new cast member"""
        self.cast.append({
            "name": "",
            "image_path": ""
        })
        self.create_cast_member_widget(len(self.cast) - 1)
    
    def remove_cast_member(self, index):
        """Remove a cast member"""
        if len(self.cast) > 1:
            # Remove from data
            self.cast.pop(index)
            
            # Remove UI elements
            for frame_data in self.cast_member_frames:
                frame_data["frame"].destroy()
            
            # Clear the list and rebuild
            self.cast_member_frames.clear()
            for i, _ in enumerate(self.cast):
                self.create_cast_member_widget(i)
        else:
            messagebox.showinfo("Cannot Remove", "You must have at least one cast member.")
    
    def update_preview(self):
        """Update the preview image (mock implementation)"""
        # In a real app, we would generate a proper preview
        # For now, just show a placeholder or the background image
        try:
            if os.path.exists(self.bg_image_path.get()):
                image = Image.open(self.bg_image_path.get())
                # Resize to fit in the preview area
                image = image.resize((400, 225), Image.LANCZOS)
                photo = ImageTk.PhotoImage(image)
                self.preview_image_label.configure(image=photo)
                self.preview_image_label.image = photo  # Keep a reference
        except Exception as e:
            self.preview_image_label.configure(text="Preview not available")
            print(f"Preview error: {e}")
    
    def update_plot_from_text(self):
        """Update the plot variable from the text widget"""
        self.plot_var.set(self.plot_text.get("1.0", tk.END).strip())
    
    def generate_powerpoint(self):
        """Generate the PowerPoint file"""
        # Update plot from text widget
        self.update_plot_from_text()
        
        # Ask user where to save the PowerPoint file
        save_filename = filedialog.asksaveasfilename(
            title="Save PowerPoint File",
            defaultextension=".pptx",
            filetypes=[("PowerPoint Presentation", "*.pptx"), ("All files", "*.*")],
            initialfile=f"{self.title_var.get().replace(' ', '_')}_Movie_Poster.pptx"
        )
        
        if not save_filename:
            # User cancelled the save dialog
            return
            
        # Ensure the filename ends with .pptx
        if not save_filename.lower().endswith('.pptx'):
            save_filename += '.pptx'
            
        # Store the filename for use in the generation process
        self.save_filename = save_filename
        
        # Simulate generation process with a progress bar
        self.status_var.set("Generating PowerPoint...")
        self.progress_var.set(0)
        
        # Generate Python code
        code = self.generate_python_code()
        
        # Save Python code if checkbox is checked
        if self.save_code_var.get():
            self.save_python_code(code)
        
        # Run the generation in a separate thread to avoid freezing UI
        threading.Thread(target=self.run_generation_process, daemon=True).start()
    
    def run_generation_process(self):
        """Run the PowerPoint generation process"""
        try:
            # Simulate progress
            for i in range(1, 101):
                self.progress_var.set(i)
                self.root.update_idletasks()
                # Add a small delay to simulate work
                import time
                time.sleep(0.02)
            
            # Here we would generate the actual PowerPoint using python-pptx
            self.create_powerpoint_file(self.save_filename)
            
            # Show success message with the path to the file
            self.status_var.set("PowerPoint generated successfully!")
            messagebox.showinfo(
                "Success", 
                f"PowerPoint has been saved to:\n{self.save_filename}"
            )
        except Exception as e:
            self.status_var.set(f"Error: {str(e)}")
            messagebox.showerror("Error", f"Failed to generate PowerPoint: {str(e)}")
            
    def create_powerpoint_file(self, filename):
        """Create the actual PowerPoint file using python-pptx"""
        try:
            # Import here to avoid startup delays
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            
            # Create presentation and add blank slide
            prs = Presentation()
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Background
            if os.path.exists(self.bg_image_path.get()):
                slide.shapes.add_picture(self.bg_image_path.get(), 0, 0, width=slide_width, height=slide_height)
            
            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
            title_frame = title_box.text_frame
            title_frame.clear()
            title_run = title_frame.paragraphs[0].add_run()
            title_run.text = self.title_var.get()
            title_run.font.size = Pt(60)
            title_run.font.bold = True
            title_run.font.name = "Impact"
            title_run.font.color.rgb = RGBColor(255, 69, 0)
            
            # Tagline
            tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(0.6))
            tagline_frame = tagline_box.text_frame
            tagline_run = tagline_frame.paragraphs[0].add_run()
            tagline_run.text = self.tagline_var.get()
            tagline_run.font.size = Pt(22)
            tagline_run.font.italic = True
            tagline_run.font.color.rgb = RGBColor(255, 255, 255)
            
            # Cast
            x_positions = [0.4, 2.0, 3.6, 5.2, 6.8][:len(self.cast)]
            if len(self.cast) > 5:
                x_positions = [n * (8 / len(self.cast)) + 0.4 for n in range(len(self.cast))]
            y_image = 2.4
            y_label = 4.0
            img_width = Inches(1.4)
            
            for i, member in enumerate(self.cast):
                x = Inches(x_positions[i])
                if os.path.exists(member["image_path"]):
                    slide.shapes.add_picture(member["image_path"], x, Inches(y_image), width=img_width)
                label_box = slide.shapes.add_textbox(x, Inches(y_label), img_width, Inches(0.5))
                label_frame = label_box.text_frame
                label_frame.paragraphs[0].text = member["name"]
                label_frame.paragraphs[0].font.size = Pt(12)
                label_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 0)
                label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Director
            director_box = slide.shapes.add_textbox(Inches(6.4), Inches(0.4), Inches(3), Inches(0.5))
            director_frame = director_box.text_frame
            director_frame.paragraphs[0].text = f"Directed by {self.director_var.get()}"
            director_frame.paragraphs[0].font.size = Pt(14)
            director_frame.paragraphs[0].font.color.rgb = RGBColor(173, 216, 230)
            
            # Plot
            plot_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(8.5), Inches(1.2))
            plot_frame = plot_box.text_frame
            plot_frame.word_wrap = True
            plot_paragraph = plot_frame.add_paragraph()
            plot_paragraph.text = self.plot_var.get()
            plot_paragraph.font.size = Pt(14)
            plot_paragraph.font.color.rgb = RGBColor(240, 240, 240)
            plot_paragraph.alignment = PP_ALIGN.LEFT
            
            # Call to Action
            cta_box = slide.shapes.add_textbox(Inches(2.5), Inches(6.8), Inches(5), Inches(1))
            cta_frame = cta_box.text_frame
            cta_run = cta_frame.paragraphs[0].add_run()
            cta_run.text = self.cta_var.get()
            cta_run.font.size = Pt(22)
            cta_run.font.bold = True
            cta_run.font.color.rgb = RGBColor(255, 0, 0)
            cta_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Save the presentation
            prs.save(filename)
            return True
        except Exception as e:
            messagebox.showerror("Error", f"Failed to create PowerPoint file: {str(e)}")
            return False
    
    def save_python_code(self, code):
        """Save the generated Python code to a file"""
        filename = filedialog.asksaveasfilename(
            defaultextension=".py",
            filetypes=[("Python files", "*.py"), ("All files", "*.*")],
            initialfile=f"{self.title_var.get().replace(' ', '_')}_generator.py"
        )
        if filename:
            with open(filename, "w") as f:
                f.write(code)
    
    def generate_python_code(self):
        """Generate the Python code for PowerPoint creation"""
        # Update plot from text widget
        self.update_plot_from_text()
        
        # Create cast code
        cast_code = "cast = [\n"
        for member in self.cast:
            cast_code += f'    ("{member["image_path"]}", "{member["name"]}"),\n'
        cast_code += "]"
        
        # Generate full code
        code = f"""from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation and add blank slide
prs = Presentation()
slide_width = prs.slide_width
slide_height = prs.slide_height
slide = prs.slides.add_slide(prs.slide_layouts[6])

# === Background ===
slide.shapes.add_picture("{self.bg_image_path.get()}", 0, 0, width=slide_width, height=slide_height)

# === Title ===
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
title_frame = title_box.text_frame
title_frame.clear()
title_run = title_frame.paragraphs[0].add_run()
title_run.text = "{self.title_var.get()}"
title_run.font.size = Pt(60)
title_run.font.bold = True
title_run.font.name = "Impact"
title_run.font.color.rgb = RGBColor(255, 69, 0)  # Red-orange

# === Tagline ===
tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(0.6))
tagline_frame = tagline_box.text_frame
tagline_run = tagline_frame.paragraphs[0].add_run()
tagline_run.text = "{self.tagline_var.get()}"
tagline_run.font.size = Pt(22)
tagline_run.font.italic = True
tagline_run.font.color.rgb = RGBColor(255, 255, 255)

# === CAST ===
{cast_code}

x_positions = [0.4, 2.0, 3.6, 5.2, 6.8][:{len(self.cast)}]
# Adjust positions based on cast count
if len(cast) > 5:
    x_positions = [n * (8 / len(cast)) + 0.4 for n in range(len(cast))]
y_image = 2.4
y_label = 4.0
img_width = Inches(1.4)

for i, (img_file, name) in enumerate(cast):
    x = Inches(x_positions[i])
    slide.shapes.add_picture(img_file, x, Inches(y_image), width=img_width)
    label_box = slide.shapes.add_textbox(x, Inches(y_label), img_width, Inches(0.5))
    label_frame = label_box.text_frame
    label_frame.paragraphs[0].text = name
    label_frame.paragraphs[0].font.size = Pt(12)
    label_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 0)
    label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# === Director ===
director_box = slide.shapes.add_textbox(Inches(6.4), Inches(0.4), Inches(3), Inches(0.5))
director_frame = director_box.text_frame
director_frame.paragraphs[0].text = "Directed by {self.director_var.get()}"
director_frame.paragraphs[0].font.size = Pt(14)
director_frame.paragraphs[0].font.color.rgb = RGBColor(173, 216, 230)

# === Plot (Moved to bottom) ===
plot_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(8.5), Inches(1.2))
plot_frame = plot_box.text_frame
plot_frame.word_wrap = True
plot_paragraph = plot_frame.add_paragraph()
plot_paragraph.text = "{self.plot_var.get()}"
plot_paragraph.font.size = Pt(14)
plot_paragraph.font.color.rgb = RGBColor(240, 240, 240)
plot_paragraph.alignment = PP_ALIGN.LEFT

# === Call to Action ===
cta_box = slide.shapes.add_textbox(Inches(2.5), Inches(6.8), Inches(5), Inches(1))
cta_frame = cta_box.text_frame
cta_run = cta_frame.paragraphs[0].add_run()
cta_run.text = "{self.cta_var.get()}"
cta_run.font.size = Pt(22)
cta_run.font.bold = True
cta_run.font.color.rgb = RGBColor(255, 0, 0)
cta_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save the final presentation
prs.save("{self.title_var.get().replace(' ', '_')}_Movie_Poster.pptx")
print("✅ Slide updated: Poster saved as {self.title_var.get().replace(' ', '_')}_Movie_Poster.pptx")
"""
        return code

def main():
    root = tk.Tk()
    app = MoviePosterApp(root)
    root.mainloop()

if __name__ == "__main__":
    main()